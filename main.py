import streamlit as st
import fitz  # PyMuPDF
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from tqdm import tqdm
from io import BytesIO
from PIL import Image

# Streamlit UI
st.title("PDF to PPTX Converter")
st.write("Upload a PDF file to convert each page to a slide in a PPTX file.")

# File uploader
uploaded_file = st.file_uploader("Choose a PDF file", type="pdf")

# Convert PDF to PPTX
def convert_pdf_to_pptx(pdf_data):
    pdf_document = fitz.open("pdf", pdf_data)
    presentation = Presentation()

    for page_num in tqdm(range(len(pdf_document)), desc="Converting PDF to PPTX"):
        page = pdf_document.load_page(page_num)
        
        # Page size in inches
        page_width = Inches(page.rect.width / 72)
        page_height = Inches(page.rect.height / 72)
        
        # Set slide size based on page size
        presentation.slide_width = page_width
        presentation.slide_height = page_height
        
        # Convert page to image
        images = convert_from_path(BytesIO(pdf_data), first_page=page_num+1, last_page=page_num+1)
        image = images[0]
        
        # Add slide and place image
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        
        image_width, image_height = image.size
        aspect_ratio = image_width / image_height

        if page_width / page_height > aspect_ratio:
            new_height = page_height
            new_width = new_height * aspect_ratio
        else:
            new_width = page_width
            new_height = new_width / aspect_ratio

        left = (page_width - new_width) / 2
        top = (page_height - new_height) / 2
        
        # Save image temporarily and add to slide
        image_bytes = BytesIO()
        image.save(image_bytes, format="PNG")
        slide.shapes.add_picture(image_bytes, left, top, width=new_width, height=new_height)
    
    # Save the presentation to a BytesIO object
    pptx_data = BytesIO()
    presentation.save(pptx_data)
    pptx_data.seek(0)
    return pptx_data

# Process PDF and provide download link
if uploaded_file is not None:
    st.write("Converting PDF to PPTX, please wait...")
    pptx_data = convert_pdf_to_pptx(uploaded_file.read())
    st.success("Conversion completed!")

    # Provide download link for PPTX
    st.download_button(
        label="Download PPTX file",
        data=pptx_data,
        file_name="converted_presentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
